from flask import Flask, request, jsonify, send_file
import os
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import time
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from langchain_core.pydantic_v1 import BaseModel, Field
from langchain_groq import ChatGroq
from langchain_community.vectorstores import Chroma
import json
import PyPDF2

from langchain_huggingface import HuggingFaceEmbeddings
from huggingface_hub import InferenceClient

from langchain import hub
from pydantic import BaseModel, Field
from typing import Dict

import fitz
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain.chains import create_retrieval_chain
from langchain_core.runnables import RunnablePassthrough
from langchain_core.tools import tool
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langgraph.checkpoint.memory import MemorySaver
from langgraph.prebuilt import ToolNode
from langchain_community.tools import DuckDuckGoSearchRun
from langchain_community.tools.tavily_search import TavilySearchResults
from langgraph.prebuilt import tools_condition
from langgraph.graph import StateGraph, MessagesState, START, END
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.pydantic_v1 import BaseModel, Field
import langchain_core
from typing import List
import langchain_community
from pprint import pprint
from IPython.display import Image, display
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain_google_genai import ChatGoogleGenerativeAI
from typing import TypedDict, List
from PIL import Image
from dotenv import load_dotenv





# from google.colab import userdata
import os
from typing import Annotated, Literal, TypedDict
import os
import time
import json
from typing import List
from langchain_core.pydantic_v1 import BaseModel, Field
# import PyPDF2
# import assemblyai as aai
from langchain.text_splitter import RecursiveCharacterTextSplitter
#@title Install requirements
from io import BytesIO
import IPython
import json
import os
from PIL import Image
import requests
import time
import os
from pptx.util import Inches  # Import Inches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# from google.colab import output
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser
from langchain_groq import ChatGroq
from langchain.schema import Document

from langchain_community.vectorstores import Chroma
from langchain_huggingface import HuggingFaceEmbeddings

from langchain import hub
import fitz  # PyMuPDF

from langchain_core.output_parsers import StrOutputParser


from PIL import Image

from langchain import hub
from langchain.chains.combine_documents import create_stuff_documents_chain

from langchain_core.prompts import ChatPromptTemplate


from typing import Dict
from langchain_core.pydantic_v1 import BaseModel, Field


from langchain.chains import create_retrieval_chain

from langchain_core.runnables import RunnablePassthrough

import re
import json

from huggingface_hub import InferenceClient


from langchain_core.tools import tool
from langchain_core.messages import HumanMessage,AIMessage,SystemMessage
from langgraph.checkpoint.memory import MemorySaver
from langgraph.prebuilt import ToolNode

from langchain_community.tools import DuckDuckGoSearchRun

from langchain_community.tools.tavily_search import TavilySearchResults
from langgraph.prebuilt import tools_condition
from langgraph.graph import StateGraph,MessagesState, START, END
from IPython.display import Image, display

from langchain_core.prompts import ChatPromptTemplate
from langchain_core.pydantic_v1 import BaseModel, Field
import langchain_core
from typing import List
import langchain_community

from pprint import pprint

from IPython.display import Image, display


from langchain_community.tools.tavily_search import TavilySearchResults

import os
from dotenv import load_dotenv
from langchain_google_genai import ChatGoogleGenerativeAI

import os
from dotenv import load_dotenv
from typing import TypedDict, List
from PIL import Image  # Import the Image class# os.environ["LANGCHAIN_API_KEY"] = LANGCHAIN_API_KEY


load_dotenv()

app = Flask(__name__)

# Constants
# GROQ_API_KEY = "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"
GROQ_API_KEY = "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"
os.environ['GROQ_API_KEY'] = GROQ_API_KEY

HUGGING_FACE_KEY = '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'
os.environ['HUGGING_FACE_KEY'] = HUGGING_FACE_KEY

# Define the State class
class State(TypedDict):
    input: str
    type: int
    section_content: List[str]
    section_name: List[int]
    section_ppt: List[str]
    prompt_for_images: List[str]
    generated_images: List[Image.Image]
    research_title: str
    professional: int

# Singleton for LLM
class GeminiSingleton:
    _instance = None

    @classmethod
    def get_instance(cls):
        if cls._instance is None:
            cls._instance = ChatGroq(model_name="llama-3.3-70b-versatile")
        return cls._instance

# SubtopicFormat class
class SubtopicFormat(BaseModel):
    section_text: Dict[int, str] = Field(description="A dictionary where keys are integers representing the section index and values are the section's text.")
    section_name: Dict[int, str] = Field(description="A dictionary where keys are integers representing the section index and values are the section's name.")

# ListFormat class
class ListFormat(BaseModel):
    points: List[str] = Field(description="A list of key points where each entry is a distinct segment of text.")

# LLM Concept Splitter
def llm_concept_splitter():
    llm = GeminiSingleton.get_instance()
    system = """
    You will be given a text and tasked with dividing it into distinct topics.
    For each topic:
    1. Extract the title and include it in a dictionary where the keys are section IDs (integers), and the values are the section names.
    2. Extract the relevant portion of the text for that topic and include it in a corresponding dictionary where the keys are section IDs (integers), and the values are the content of each section.
    Guidelines:
    - Ensure that each topic includes only relevant concepts or ideas from the context. Avoid adding meta-text like "This section deals with XYZ."
    - Maintain a one-to-one mapping between the section IDs, section names, and section contents.
    Your output should strictly be a JSON string containing exactly two keys:
    1. "section_text": A dictionary where each key is a section ID (integer), and the value is the corresponding section content.
    2. "section_name": A dictionary where each key is a section ID (integer), and the value is the name of the corresponding section.
    Do not include anything else, not even ```json, only json content without anything else. Make sure to give correct json response by enclosing both the key and value in double quotes.
    """
    concept_split_prompt = ChatPromptTemplate.from_messages([("system", system), ("human", "Here is the text:{text}.")])
    llm_with_structured_output = llm.with_structured_output(SubtopicFormat)
    # return concept_split_prompt | llm_with_structured_output
    return concept_split_prompt | llm

# Extract text from PDF
# def extract_text_from_pdf(pdf_path: str) -> str:
#     doc = fitz.open(pdf_path)
#     text = ""
#     for page in doc:
#         text += page.get_text("text") + "\n\n"
#     return text

import fitz  # PyMuPDF

# def extract_text_from_pdf(pdf_file):
#     try:
#         pdf_data = BytesIO(pdf_file.read())  # Convert the uploaded PDF to a file-like object
        
#         # Log the length of the data for debugging
#         print(f"PDF data length: {len(pdf_data.getvalue())}")
        
#         # Open the PDF with fitz.open
#         pdf_document = fitz.open(pdf_data)
        
#         text = ""
#         for page_num in range(pdf_document.page_count):
#             page = pdf_document.load_page(page_num)
#             text += page.get_text()
        
#         return text
#     except Exception as e:
#         print(f"Error extracting text from PDF: {e}")
#         return None


import PyPDF2
from io import BytesIO

def extract_text_from_pdf(pdf_file):
    """Extracts full text from a PDF file-like object."""
    text = ""
    try:
        # Convert the FileStorage object to a file-like object
        pdf_data = BytesIO(pdf_file.read())
        
        # Open the PDF with PyPDF2
        reader = PyPDF2.PdfReader(pdf_data)
        
        for page in reader.pages:
            text += page.extract_text() + "\n"
        
        return text.strip()
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
        return None



# PowerPoint Presenter
def ppt_presentor(text):
    llm = GeminiSingleton.get_instance()
    system = """
    You will be given a text and tasked with summarizing it concisely so that it can be effectively presented on a single PowerPoint slide.
    Instructions:
    1. Extract the key ideas and main points from the provided text.
    2. Ensure the summary is brief, clear, and suitable for a single slide.
    3. Do not use bullet points, numbering, or special characters—just plain text.
    4. Keep the language simple and professional, ensuring readability.
    5. Maintain coherence while removing redundant details and complex explanations.
    """
    msg = f"{system}.{text}"
    result = llm.invoke(msg)
    return result.content

# Prompt Generator
def prompt_generator(text):
    llm = GeminiSingleton.get_instance()
    system = """
    You will be given a text and tasked with converting it into a format that can be used by an image generation model (Flux model).
    Instructions:
    1. Extract the key visual concepts and main ideas from the provided text.
    2. Transform these concepts into a detailed, descriptive format suitable for image generation.
    3. Focus on visual characteristics such as colors, objects, environments, characters, and actions that can be represented visually.
    4. Avoid using bullet points, numbering, or special characters—just plain text.
    5. Ensure that the description is clear and specific to guide the image generator in creating a detailed and accurate image.
    """
    msg = f"{system}. The text is:{text}"
    result = llm.invoke(msg)
    return result.content

# Image Generation
def generate_image(input_text: str):
    """
    Generates an image using the FLUX model based on the input text.

    Args:
        input_text (str): The text prompt for image generation.
        api_key (str): The Hugging Face API key for authentication.

    Returns:
        PIL.Image.Image: The generated image.
    """

    # print(HUGGING_FACE_KEY)
    client = InferenceClient(
        # provider="replicate",
        api_key=HUGGING_FACE_KEY
    )

    # Generate image based on the input text
    image = client.text_to_image(
        input_text,
        model="black-forest-labs/FLUX.1-dev"
    )

    # image_path = "output.png"
    # image = Image.open(image_path)

    return image



def llm_point_splitter():
    # llm = lb.LLMSingleton.get_instance()
    llm = GeminiSingleton.get_instance()

    # Define the system prompt for rewriting the question
    
    system = f"""You will be given a text and tasked with dividing it into distinct topics.
For each topic:
1. Extract the title and include it in a dictionary where the keys are section IDs (integers), and the values are the section names.
2. Extract the relevant portion of the text for that topic and include it in a corresponding dictionary where the keys are section IDs (integers), and the values are the content of each section.

Guidelines:
- Ensure that each topic includes only relevant concepts or ideas from the context. Avoid adding meta-text like "This section deals with XYZ."
- Maintain a one-to-one mapping between the section IDs, section names, and section contents. The section ID at index 0 should correspond to the content and name at index 0, and so on.
- Only Create at maximum 3 points not more than that.
Your output should strictly be a JSON string containing exactly two keys:
1. `"section_text"`: A dictionary where each key is a section ID (integer), and the value is the corresponding section content.
2. `"section_name"`: A dictionary where each key is a section ID (integer), and the value is the name of the corresponding section.
Do not include any other text like ```json. Only give the json directly


Example output:
"\"points\": ['', '', '']
"

where
points: List[str] = Field(
        description="A list of key points where each entry is a distinct segment of text."
    )

"""
    # Define the rewrite prompt template
    concept_split_prompt = ChatPromptTemplate.from_messages(
        [
            ("system", system),
            (
                "human",
                "Here is the text:{text}."
            ),
        ]
    )

    llm_with_structured_output = llm.with_structured_output(ListFormat)
    # llm_with_structured_output = client.with_structured_output(SubtopicFormat)



    # Create the question rewriter with the defined prompt and output parser
    # concept_splitter = concept_split_prompt | llm | StrOutputParser()
    # concept_splitter = concept_split_prompt | llm_with_structured_output | StrOutputParser()


    # concept_splitter = concept_split_prompt | llm

    # llm_with_structured_output.invoke({"text": text})

    # Return the question rewriter
    # return concept_splitter
    return llm_with_structured_output



# Chatbot Question Generator
class Chatbot_Question_Generator:
    def __init__(self):
        self.llm = GeminiSingleton.get_instance()
        self.llm_for_splitter = llm_concept_splitter()
        self.llm_for_point = llm_point_splitter()

    def node1(self, state):
        input = state["input"]
        old_section_names = []
        old_section_texts = []
        extracted_text = input
        chunk_size = 3000

        for i in range(0, len(extracted_text), chunk_size):
            chunk = extracted_text[i:i + chunk_size]
            reply = self.llm_for_splitter.invoke({"text": chunk})
            json_str = reply.content.replace("'", '"')
            data = json.loads(json_str)
            old_section_names.append(data['section_name'])
            old_section_texts.append(data['section_text'])

        section_texts = []
        section_names = []

        for section_dict in old_section_names:
            for key, value in section_dict.items():
                section_names.append(value)

        for section_dict in old_section_texts:
            for key, value in section_dict.items():
                section_texts.append(value)

        return {"section_name": section_names, "section_content": section_texts}

    def node2(self, state):
        section_texts = state["section_content"]
        section_ppt = []

        for section_text in section_texts:
            text = section_text
            store = ppt_presentor(text[:3000])
            section_ppt.append(store)
            time.sleep(15)

        return {"section_ppt": section_ppt}

    def node3(self, state):
        section_ppt = state["section_ppt"]
        prompt_for_image = []

        for ppt in section_ppt:
            prompt = prompt_generator(ppt)
            prompt_for_image.append(prompt[:3000])
            time.sleep(15)

        return {"prompt_for_images": prompt_for_image}

    def node4(self, state):
        prompt_for_image = state["prompt_for_images"]
        generated_images = []

        for prompt in prompt_for_image:
            image = generate_image(prompt)
            generated_images.append(image)

        return {"generated_images": generated_images}

    # def node5(self, state):
    #     try:
    #         section_content = state['section_content']
    #         generated_images = state['generated_images']
    #         section_name = state['section_name']

    #         prs = Presentation()

    #         TITLE_COLOR = RGBColor(255, 255, 255)
    #         BANNER_COLOR = RGBColor(0, 102, 204)
    #         TEXT_COLOR = RGBColor(0, 0, 0)
    #         CONTENT_BG_COLOR = RGBColor(230, 230, 230)

    #         RESEARCH_TITLE = state["research_title"]
    #         first_slide = prs.slides.add_slide(prs.slide_layouts[5])

    #         title_banner = first_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1))
    #         title_banner.fill.solid()
    #         title_banner.fill.fore_color.rgb = BANNER_COLOR
    #         title_banner.line.color.rgb = RGBColor(255, 255, 255)

    #         title_shape = first_slide.shapes.add_textbox(Inches(0.3), Inches(0.2), prs.slide_width - Inches(0.6), Inches(0.8))
    #         title_frame = title_shape.text_frame
    #         title_frame.text = RESEARCH_TITLE
    #         title_paragraph = title_frame.paragraphs[0]
    #         title_paragraph.font.bold = True
    #         title_paragraph.font.size = Pt(32)
    #         title_paragraph.font.color.rgb = TITLE_COLOR
    #         title_paragraph.alignment = PP_ALIGN.CENTER

    #         toc_shape = first_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
    #         toc_frame = toc_shape.text_frame
    #         toc_frame.text = "Table of Contents"
    #         toc_paragraph = toc_frame.paragraphs[0]
    #         toc_paragraph.font.size = Pt(28)
    #         toc_paragraph.font.bold = True
    #         toc_paragraph.alignment = PP_ALIGN.LEFT

    #         table_of_contents = section_name

    #         for section in table_of_contents:
    #             toc_item = toc_frame.add_paragraph()
    #             toc_item.text = section
    #             toc_item.font.size = Pt(20)
    #             toc_item.alignment = PP_ALIGN.LEFT

    #         is_image_left = True

    #         for title, section_text, generated_image in zip(section_name, section_content, generated_images):
    #             generated_image.save("output.png", format="PNG")
    #             bullet_points = self.llm_for_point.invoke(section_text[:3000]).points

    #             slide = prs.slides.add_slide(prs.slide_layouts[5])
    #             slide_width = prs.slide_width
    #             slide_height = prs.slide_height

    #             title_banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), slide_width, Inches(1))
    #             title_banner.fill.solid()
    #             title_banner.fill.fore_color.rgb = BANNER_COLOR
    #             title_banner.line.color.rgb = RGBColor(255, 255, 255)

    #             title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), slide_width - Inches(0.6), Inches(0.8))
    #             title_frame = title_shape.text_frame
    #             title_frame.text = title
    #             title_paragraph = title_frame.paragraphs[0]
    #             title_paragraph.font.bold = True
    #             title_paragraph.font.size = Pt(32)
    #             title_paragraph.font.color.rgb = TITLE_COLOR
    #             title_paragraph.alignment = PP_ALIGN.CENTER

    #             content_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), slide_width - Inches(1), slide_height - Inches(2.0))
    #             content_bg.fill.solid()
    #             content_bg.fill.fore_color.rgb = CONTENT_BG_COLOR
    #             content_bg.line.color.rgb = RGBColor(180, 180, 180)

    #             text_box_width = slide_width - Inches(5.1)
    #             text_box_height = slide_height - Inches(2.5)
    #             text_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5) if len(prs.slides) % 2 != 0 else slide_width - text_box_width - Inches(0.5), Inches(1.8), text_box_width, text_box_height)
    #             text_box.fill.solid()
    #             text_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    #             text_box.line.color.rgb = RGBColor(0, 51, 102)

    #             text_frame = text_box.text_frame
    #             text_frame.word_wrap = True
    #             text_frame.margin_top = Inches(0.3)
    #             text_frame.margin_bottom = Inches(0.3)
    #             text_frame.margin_left = Inches(0.3)
    #             text_frame.margin_right = Inches(0.3)

    #             for i, point in enumerate(bullet_points[:3], start=1):
    #                 p = text_frame.add_paragraph()
    #                 p.text = f"{i}. {point}"
    #                 p.font.size = Pt(18)
    #                 p.font.bold = False
    #                 p.font.color.rgb = TEXT_COLOR
    #                 p.alignment = PP_ALIGN.LEFT
    #                 p.numbered = True
    #                 p.space_after = Pt(6)

    #             try:
    #                 image_left = Inches(0.5) if len(prs.slides) % 2 == 0 else slide_width / 2 + Inches(0.5)
    #                 image_top = Inches(3.2)
    #                 img = slide.shapes.add_picture("output.png", image_left, image_top, width=Inches(3.5))
    #                 img.shadow.inherit = False
    #             except Exception as e:
    #                 print(f"Error adding image: {e}")

    #             is_image_left = not is_image_left

    #         last_slide = prs.slides.add_slide(prs.slide_layouts[5])

    #         title_banner = last_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1))
    #         title_banner.fill.solid()
    #         title_banner.fill.fore_color.rgb = BANNER_COLOR
    #         title_banner.line.color.rgb = RGBColor(255, 255, 255)

    #         thank_you_shape = last_slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    #         thank_you_frame = thank_you_shape.text_frame
    #         thank_you_frame.text = "Thank You!\nFeel free to ask any questions."
    #         thank_you_paragraph = thank_you_frame.paragraphs[0]
    #         thank_you_paragraph.font.size = Pt(36)
    #         thank_you_paragraph.font.bold = True
    #         thank_you_paragraph.font.color.rgb = TEXT_COLOR
    #         thank_you_paragraph.alignment = PP_ALIGN.CENTER

    #         ppt_file = 'new_ppt.pptx'
    #         if os.path.exists(ppt_file):
    #             os.remove(ppt_file)
    #         prs.save("new_ppt.pptx")
    #         time.sleep(15)
    #         print("PowerPoint presentation created successfully!")

    #     except Exception as e:
    #         print(f"An error occurred: {e}")
    #         prs.save("new_ppt.pptx")
    def node5(self, state):
        try:
            section_content = state['section_content']
            generated_images = state['generated_images']
            section_name = state['section_name']

            prs = Presentation()
            if state["professional"] == '1':

                # name_dict = section_name[-1]
                # text_dict = section_content[-1]
                TITLE_COLOR = RGBColor(255, 255, 255)  # White
                BANNER_COLOR = RGBColor(0, 102, 204)  # Dark Blue
                TEXT_COLOR = RGBColor(0, 0, 0)  # Black
                CONTENT_BG_COLOR = RGBColor(230, 230, 230)  # Light Gray




                # First Slide: Research Paper Title & Table of Contents
                RESEARCH_TITLE = state["research_title"]
                first_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

                # Title Banner
                title_banner = first_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1)
                )
                title_banner.fill.solid()
                title_banner.fill.fore_color.rgb = BANNER_COLOR
                title_banner.line.color.rgb = RGBColor(255, 255, 255)  # White border


                # Research Title
                title_shape = first_slide.shapes.add_textbox(Inches(0.3), Inches(0.2), prs.slide_width - Inches(0.6), Inches(0.8))
                title_frame = title_shape.text_frame
                title_frame.text = RESEARCH_TITLE
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.bold = True
                title_paragraph.font.size = Pt(32)
                title_paragraph.font.color.rgb = TITLE_COLOR
                title_paragraph.alignment = PP_ALIGN.CENTER


                toc_shape = first_slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
                toc_frame = toc_shape.text_frame
                toc_frame.text = "Table of Contents"
                toc_paragraph = toc_frame.paragraphs[0]
                toc_paragraph.font.size = Pt(28)
                toc_paragraph.font.bold = True
                toc_paragraph.alignment = PP_ALIGN.LEFT


                # List of sections (Generated from slide titles)
                table_of_contents = section_name

                # Add numbered Table of Contents
                for section in table_of_contents:
                    toc_item = toc_frame.add_paragraph()
                    toc_item.text = section
                    toc_item.font.size = Pt(20)
                    toc_item.alignment = PP_ALIGN.LEFT





                is_image_left = True

                # for (section_id, title), (section_text_id, section_text), generated_image in zip(name_dict.items(), text_dict.items(), generated_images):
                # for (section_id, title), (section_text_id, section_text), generated_image in zip(name_dict.items(), text_dict.items(), generated_images):
                for title, section_text, generated_image in zip(section_name, section_content, generated_images):

                    generated_image.save("output.png", format="PNG")
                    print("Image saved successfully!")

                    bullet_points = self.llm_for_point.invoke(section_text[:3000]).points
                    
                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height

                    title_banner = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), slide_width, Inches(1)
                    )
                    title_banner.fill.solid()
                    title_banner.fill.fore_color.rgb = BANNER_COLOR
                    title_banner.line.color.rgb = RGBColor(255, 255, 255)  # White border

                    title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), slide_width - Inches(0.6), Inches(0.8))
                    title_frame = title_shape.text_frame
                    title_frame.text = title
                    title_paragraph = title_frame.paragraphs[0]
                    title_paragraph.font.bold = True
                    title_paragraph.font.size = Pt(32)
                    title_paragraph.font.color.rgb = TITLE_COLOR
                    title_paragraph.alignment = PP_ALIGN.CENTER

                    content_bg = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), slide_width - Inches(1), slide_height - Inches(2.0)
                    )
                    content_bg.fill.solid()
                    content_bg.fill.fore_color.rgb = CONTENT_BG_COLOR
                    content_bg.line.color.rgb = RGBColor(180, 180, 180)

                    text_box_width = slide_width - Inches(5.1)
                    text_box_height = slide_height - Inches(2.5)
                    text_box = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.5) if len(prs.slides) % 2 != 0 else slide_width - text_box_width - Inches(0.5),
                        Inches(1.8),
                        text_box_width,
                        text_box_height
                    )
                    text_box.fill.solid()
                    text_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    text_box.line.color.rgb = RGBColor(0, 51, 102)

                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True
                    text_frame.margin_top = Inches(0.3)
                    text_frame.margin_bottom = Inches(0.3)
                    text_frame.margin_left = Inches(0.3)
                    text_frame.margin_right = Inches(0.3)

                    for i, point in enumerate(bullet_points[:3], start=1):
                        p = text_frame.add_paragraph()
                        p.text = f"{i}. {point}"
                        p.font.size = Pt(18)
                        p.font.bold = False
                        p.font.color.rgb = TEXT_COLOR
                        p.alignment = PP_ALIGN.LEFT
                        p.numbered = True
                        p.space_after = Pt(6)

                    try:
                        image_left = Inches(0.5) if len(prs.slides) % 2 == 0 else slide_width / 2 + Inches(0.5)
                        image_top = Inches(3.2)
                        img = slide.shapes.add_picture("output.png", image_left, image_top, width=Inches(3.5))
                        img.shadow.inherit = False
                    except Exception as e:
                        print(f"Error adding image: {e}")

                    is_image_left = not is_image_left

                
                # Last Slide: Thank You Message
                last_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

                # Title Banner
                title_banner = last_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1)
                )
                title_banner.fill.solid()
                title_banner.fill.fore_color.rgb = BANNER_COLOR
                title_banner.line.color.rgb = RGBColor(255, 255, 255)

                
                # Thank You Message
                thank_you_shape = last_slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
                thank_you_frame = thank_you_shape.text_frame
                thank_you_frame.text = "Thank You!\nFeel free to ask any questions."
                thank_you_paragraph = thank_you_frame.paragraphs[0]
                thank_you_paragraph.font.size = Pt(36)
                thank_you_paragraph.font.bold = True
                thank_you_paragraph.font.color.rgb = TEXT_COLOR
                thank_you_paragraph.alignment = PP_ALIGN.CENTER


                ppt_file = 'new_ppt.pptx'
                if os.path.exists(ppt_file):
                    os.remove(ppt_file)
                    print(f"File {ppt_file} has been deleted.")
                else:
                    print(f"File {ppt_file} not found. No action taken.")

                prs.save("new_ppt.pptx")
                # time.sleep(7)
                print("PowerPoint presentation created successfully!")
            else:

                # Create a presentation object
                

                # Define fun-themed colors
                TITLE_COLOR = RGBColor(255, 255, 255)  # White
                BANNER_COLOR = RGBColor(255, 165, 0)   # Fun Orange
                TEXT_COLOR = RGBColor(50, 50, 50)      # Dark Gray
                CONTENT_BG_COLOR = RGBColor(255, 255, 204)  # Soft Yellow
                ACCENT_COLOR = RGBColor(102, 204, 255)  # Sky Blue

                # Research Paper Title and Table of Contents
                research_title = state["research_title"]
                table_of_contents = section_name

                # Add first slide (Title & Table of Contents)
                first_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

                # Add title banner
                banner = first_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
                )
                banner.fill.solid()
                banner.fill.fore_color.rgb = BANNER_COLOR

                # Title text
                title_box = first_slide.shapes.add_textbox(Inches(0.3), Inches(0.3), prs.slide_width - Inches(0.6), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = research_title
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.size = Pt(34)
                title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Table of Contents
                content_box = first_slide.shapes.add_textbox(Inches(1), Inches(1.8), prs.slide_width - Inches(2), Inches(3))
                content_frame = content_box.text_frame
                content_frame.text = "Table of Contents"
                content_frame.paragraphs[0].font.bold = True
                content_frame.paragraphs[0].font.size = Pt(28)
                content_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
                content_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

                for section in table_of_contents:
                    p = content_frame.add_paragraph()
                    p.text = section
                    p.font.size = Pt(22)
                    p.font.color.rgb = TEXT_COLOR
                    p.alignment = PP_ALIGN.LEFT
                    p.space_after = Pt(8)

                # Sample slides with titles, bullet points, and image paths
                

                # for title, bullet_points, image_path in slides_content:
                for title, section_text, generated_image in zip(section_name, section_content, generated_images):
                    
                    bullet_points = self.llm_for_point.invoke(section_text[:3000]).points
                    generated_image.save("output.png", format="PNG")
                    print("Image saved successfully!")

                    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height

                    # Add title banner
                    title_banner = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0), Inches(0), slide_width, Inches(1.2)
                    )
                    title_banner.fill.solid()
                    title_banner.fill.fore_color.rgb = BANNER_COLOR

                    # Add title text
                    title_shape = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), slide_width - Inches(0.6), Inches(1))
                    title_frame = title_shape.text_frame
                    title_frame.text = title
                    title_frame.paragraphs[0].font.bold = True
                    title_frame.paragraphs[0].font.size = Pt(34)
                    title_frame.paragraphs[0].font.color.rgb = TITLE_COLOR
                    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

                    # Add content background (text on left)
                    text_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), slide_width/2 - Inches(1), slide_height - Inches(2.0)
                    )
                    text_bg.fill.solid()
                    text_bg.fill.fore_color.rgb = CONTENT_BG_COLOR

                    # Add image background (right side)
                    image_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, slide_width/2, Inches(1.5), slide_width/2 - Inches(0.5), slide_height - Inches(2.0)
                    )
                    image_bg.fill.solid()
                    image_bg.fill.fore_color.rgb = RGBColor(240, 240, 240)

                    # Add text box
                    text_box = slide.shapes.add_textbox(Inches(1), Inches(2), slide_width/2 - Inches(2), slide_height - Inches(3))
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = True

                    for i, point in enumerate(bullet_points, start=1):
                        if i > 3:
                            break
                        p = text_frame.add_paragraph()
                        p.text = f"{i}. {point}"
                        p.font.size = Pt(20)
                        p.font.bold = True
                        p.font.color.rgb = TEXT_COLOR

                    try:
                        image_left = slide_width / 2 + Inches(0.6)
                        image_top = Inches(3.2)
                        img = slide.shapes.add_picture("output.png", image_left, image_top, width=Inches(3.5))
                        img.shadow.inherit = False
                    except Exception as e:
                        print(f"Error adding image: {e}")

                    # is_image_left = not is_image_left

                    # Add image
                    # try:
                    #     slide.shapes.add_picture(image_path, slide_width/2 + Inches(0.5), Inches(2.5), width=slide_width/2 - Inches(2))
                    # except Exception as e:
                    #     print(f"Error adding image {image_path}: {e}")

                # Add last slide (Thank You)
                last_slide = prs.slides.add_slide(prs.slide_layouts[5])

                title_box = last_slide.shapes.add_textbox(Inches(1), Inches(3), prs.slide_width - Inches(2), Inches(3))
                title_frame = title_box.text_frame
                title_frame.text = "🙏 Thank You for Your Attention! 🎈"
                title_frame.paragraphs[0].font.bold = True
                title_frame.paragraphs[0].font.size = Pt(40)
                title_frame.paragraphs[0].font.color.rgb = BANNER_COLOR
                title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER



                ppt_file = 'new_ppt.pptx'
                if os.path.exists(ppt_file):
                    os.remove(ppt_file)
                    print(f"File {ppt_file} has been deleted.")
                else:
                    print(f"File {ppt_file} not found. No action taken.")

                prs.save("new_ppt.pptx")
                # time.sleep(7)
                print("PowerPoint presentation created successfully!")

        except Exception as e:
            print(f"An error occurred: {e}")
            prs.save("new_ppt.pptx")


    def __call__(self):
        workflow = StateGraph(State)
        workflow.add_node("node1", self.node1)
        workflow.add_node("node2", self.node2)
        workflow.add_node("node3", self.node3)
        workflow.add_node("node4", self.node4)
        workflow.add_node("node5", self.node5)

        workflow.add_edge(START, "node1")
        workflow.add_edge("node1", "node2")
        workflow.add_edge("node2", "node3")
        workflow.add_edge("node3", "node4")
        workflow.add_edge("node4", "node5")
        workflow.add_edge("node5", END)

        self.app = workflow.compile()
        return self.app

# Initialize the chatbot
mybot = Chatbot_Question_Generator()
workflow = mybot()

# Flask endpoint
# @app.route('/generate_ppt', methods=['POST'])
# def generate_ppt():
#     try:
#         data = request.json
#         input_text = data.get('input')
#         research_title = data.get('research_title')

#         if not input_text or not research_title:
#             return jsonify({"error": "Missing 'input' or 'research_title' in request"}), 400

#         response = workflow.invoke({"input": input_text, "research_title": research_title})
#         return send_file("new_ppt.pptx", as_attachment=True)

#     except Exception as e:
#         return jsonify({"error": str(e)}), 500



@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    try:
        if 'pdf_file' not in request.files or 'research_title' not in request.form:
            return jsonify({"error": "Missing 'pdf_file' or 'research_title' in request"}), 400
        
        pdf_file = request.files['pdf_file']
        print("type of pdf.......",type(pdf_file))
        research_title = request.form['research_title']

        professional = request.form['professional']

        # Ensure it's a PDF
        if not pdf_file.filename.endswith('.pdf'):
            return jsonify({"error": "Invalid file format. Only PDFs are allowed."}), 400

        # Extract text from the uploaded PDF
        pdf_text = extract_text_from_pdf(pdf_file)
        print("here...")
        # return jsonify({"Success": pdf_text}), 200
        # Invoke the workflow with extracted text
        extract = pdf_text[:1500]
        # print("The professional is...", professional)
        # print("type", type(professional))
        # return 
        response = workflow.invoke({"input": extract, "research_title": research_title, "professional": professional})
        # response = workflow.invoke({"input": pdf_text, "research_title": research_title, "professional": professional})

        # Return the generated PowerPoint file
        return send_file("new_ppt.pptx", as_attachment=True)

    except Exception as e:
        return jsonify({"error": str(e)}), 500

# if __name__ == '__main__':
    # app.run(debug=True)

if __name__ == '__main__':
    app.run(debug=True, port=8080)
