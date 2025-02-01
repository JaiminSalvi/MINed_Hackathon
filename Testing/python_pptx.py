from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Create a new PowerPoint presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

# Set Gradient Background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(20, 20, 40)  # Deep blue gradient

# Add a Stylish Shape as a Header
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(9), Inches(1))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold color
shape.shadow.inherit = False

# Add Title with Bold Gold Text
title_shape = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8), Inches(1))
title_text_frame = title_shape.text_frame
title_text_frame.text = "Introduction to BERT"
title_text_frame.paragraphs[0].font.size = Pt(40)
title_text_frame.paragraphs[0].font.bold = True
title_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text for contrast
title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add Subtitle with Elegant Styling
subtitle_shape = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(0.5))
subtitle_text_frame = subtitle_shape.text_frame
subtitle_text_frame.text = "Bidirectional Encoder Representations from Transformers"
subtitle_text_frame.paragraphs[0].font.size = Pt(24)
subtitle_text_frame.paragraphs[0].font.color.rgb = RGBColor(173, 216, 230)  # Light Blue
subtitle_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add a Callout Shape for Key Features
content_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(6), Inches(3.5))
content_shape.fill.solid()
content_shape.fill.fore_color.rgb = RGBColor(40, 40, 60)  # Dark blue
content_shape.shadow.inherit = False

# Add Key Points in Bullet Format Inside the Callout
content_text_frame = content_shape.text_frame
content_text_frame.word_wrap = True
content_text_frame.text = "🔹 Key Features of BERT:"

# List of key points with icons
key_points = [
    "📌 Transformer-based model by Google",
    "📌 Deeply bidirectional for better context understanding",
    "📌 Pre-trained with Masked Language Modeling",
    "📌 State-of-the-art in NLP tasks like QA",
    "📌 Used in search engines, chatbots, AI assistants"
]

for point in key_points:
    para = content_text_frame.add_paragraph()
    para.text = point
    para.font.size = Pt(18)
    para.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add a High-Quality Image of BERT Model
img_path = "image.png"  # Ensure the image file exists
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(7), Inches(2.5), width=Inches(3.5))
else:
    print("Warning: Image not found. Please check the path.")

# Add a Footer for Extra Style
footer_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.5), Inches(10), Inches(0.5))
footer_shape.fill.solid()
footer_shape.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold footer

# Save the Presentation
prs.save("BERT_Styled_Intro_Slide.pptx")

print("Stylish BERT Introduction Slide Created Successfully!")
